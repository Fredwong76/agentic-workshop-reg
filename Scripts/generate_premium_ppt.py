from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_text_style(text_frame, text, font_size=18, color=RGBColor(255, 255, 255), bold=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    return p

def create_premium_slide():
    prs = Presentation()
    
    # Set Slide Size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Blank Layout
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    
    # 1. Background (Dark Blue/Slate)
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height # 1 is Rectangle
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(15, 23, 42) # slate-900
    background.line.fill.background()

    # 2. Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.text = "Gap Analysis & Interview Ready"
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(56, 189, 248) # light blue accent
    
    # 3. Card 1: Match
    card1_bg = slide.shapes.add_shape(
        5, Inches(0.5), Inches(1.8), Inches(6), Inches(2.5) # 5 is Rounded Rectangle
    )
    card1_bg.fill.solid()
    card1_bg.fill.fore_color.rgb = RGBColor(30, 41, 59) # slate-800
    card1_bg.line.color.rgb = RGBColor(71, 85, 105) # slate-600
    
    txt1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(2))
    tf1 = txt1.text_frame
    tf1.word_wrap = True
    apply_text_style(tf1, "💪 Strongest Matches", 24, RGBColor(56, 189, 248), True)
    apply_text_style(tf1, "• 90% technical match for Azure/AWS migrations.", 18)
    apply_text_style(tf1, "• Deep enterprise assessment experience.", 18)

    # 4. Card 2: Gap
    card2_bg = slide.shapes.add_shape(
        5, Inches(6.8), Inches(1.8), Inches(6), Inches(2.5)
    )
    card2_bg.fill.solid()
    card2_bg.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card2_bg.line.color.rgb = RGBColor(239, 68, 68) # red accent
    
    txt2 = slide.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.5), Inches(2))
    tf2 = txt2.text_frame
    tf2.word_wrap = True
    apply_text_style(tf2, "⚠️ Identified Gaps", 24, RGBColor(239, 68, 68), True)
    apply_text_style(tf2, "• No mention of GovZTA (Zero Trust Architecture).", 18)
    apply_text_style(tf2, "• Missing specific Singapore Govt compliance history.", 18)

    # 5. Bottom Banner (Strategic Question)
    banner = slide.shapes.add_shape(
        1, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(56, 189, 248)
    banner.fill.transparency = 0.9 # Subtle glow
    banner.line.color.rgb = RGBColor(56, 189, 248)
    
    txt3 = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11), Inches(1.5))
    tf3 = txt3.text_frame
    tf3.word_wrap = True
    apply_text_style(tf3, "🎯 Recommended Killer Question:", 22, RGBColor(56, 189, 248), True)
    apply_text_style(tf3, "\"Explain how you would implement Zero Trust principles in a highly regulated government environment vs. a standard enterprise deployment.\"", 20, RGBColor(255, 255, 255))

    # Save
    output_path = r"Workshop_Materials\Premium_Python_Presentation.pptx"
    prs.save(output_path)
    print(f"Premium Presentation saved to {output_path}")

if __name__ == "__main__":
    create_premium_slide()
