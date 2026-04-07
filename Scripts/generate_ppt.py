from pptx import Presentation
from pptx.util import Inches, Pt

def create_slide():
    prs = Presentation()
    
    # Title Slide
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Gap Analysis & Interview Ready"
    
    # Setting up the content
    tf = content.text_frame
    tf.text = "Project: Interviews (Eric Ho)"
    
    p = tf.add_paragraph()
    p.text = "Status: 30-second strategic analysis complete."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Key Findings:"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Match: 90% Azure/AWS technical requirements."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Gap: GovZTA (Government Zero Trust Architecture) missing."
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "Strategic Question:"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "'Explain how you would implement Zero Trust principles in a highly regulated environment.'"
    p.level = 2
    
    # Save the presentation
    output_path = r"Workshop_Materials\PanDoc_Alternative_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_slide()
